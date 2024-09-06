# Copyright (C) 2024 Intel Corporation
# SPDX-License-Identifier: Apache-2.0

import errno
import functools
import io
import json
import os
import re
import shutil
import signal
import time
import unicodedata
from pathlib import Path
from typing import List
import uuid

import docx
import docx2txt
import easyocr
from fastapi import UploadFile
import fitz
import numpy as np
import pandas as pd
import pptx
import yaml
from langchain_community.document_loaders import (
    UnstructuredHTMLLoader,
    UnstructuredMarkdownLoader,
    UnstructuredXMLLoader,
)
from PIL import Image

from comps.cores.proto.docarray import TextDoc
from comps.dataprep.utils.crawler import Crawler
import logging

class TimeoutError(Exception):
    pass

async def save_file_to_local_disk(file: UploadFile) -> str:
    upload_folder =  os.getenv('UPLOAD_PATH', '/tmp/opea_upload')

    if not os.path.exists(upload_folder):
        Path(upload_folder).mkdir(parents=True, exist_ok=True)

    save_path = Path(os.path.join(upload_folder, file.filename))
    with save_path.open("wb") as fout:
        try:
            content = await file.read()
            fout.write(content)
            fout.close
        except Exception as e:
            logging.exception(f"Write file failed. Exception: {e}")
            raise

    return save_path

async def save_link_to_local_disk(link_list: List[str]) -> List[str]:
    upload_folder =  os.getenv('UPLOAD_PATH', '/tmp/opea_upload')
    if not os.path.exists(upload_folder):
        Path(upload_folder).mkdir(parents=True, exist_ok=True)

    data_collection = parse_html(link_list) # crawl through the page

    save_paths = []
    for data, meta in data_collection:
        doc_id = str(uuid.uuid4())+".txt"
        save_path = Path(os.path.join(upload_folder, doc_id))
        with save_path.open("w") as fout:
            try:
                fout.write(data)
                fout.close
                save_paths.append(save_path)
            except Exception as e:
                logging.exception(f"Write file failed. Exception: {e}")
                raise
    return save_paths

def timeout(seconds=10, error_message=os.strerror(errno.ETIME)):
    def decorator(func):
        def _handle_timeout(signum, frame):
            raise TimeoutError(error_message)

        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            signal.signal(signal.SIGALRM, _handle_timeout)
            signal.alarm(seconds)
            try:
                result = func(*args, **kwargs)
            finally:
                signal.alarm(0)
            return result

        return wrapper

    return decorator

async def parse_files(files: List[UploadFile]) -> List[TextDoc]:
    parsed_texts: List[TextDoc] = []

    for file in files:
        # if files have to be persisted internally
        try:
            path = await save_file_to_local_disk(file)
            saved_path = str(path.resolve())
            logging.info(f"saved file {file.filename} to {saved_path}")

            metadata = {
                'path': saved_path,
                'timestamp': time.time()
            }

            parsed_texts.append(
                TextDoc(
                    text=document_loader(saved_path),
                    metadata = metadata
                )
            )
        except Exception as e:
            logging.exception(e)
            raise e

    return parsed_texts

async def parse_links(links: List[str]) -> List[TextDoc]:
    parsed_texts: List[TextDoc] = []

    for link in links:
        try:
            paths = await save_link_to_local_disk([link])
            for path in paths:
                saved_path = str(path.resolve())
                logging.info(f"saved link {link} to {saved_path}")

                metadata = {
                    'path': saved_path,
                    'url': link,
                    'timestamp': time.time()
                }

                parsed_texts.append(
                    TextDoc(
                        text=document_loader(saved_path),
                        metadata = metadata
                    )
                )

        except Exception as e:
            logging.exception(e)
            raise e

    return parsed_texts

def get_split_headers():
    return [
        ("h1", "Header 1"),
        ("h2", "Header 2"),
        ("h3", "Header 3"),
    ]

def get_separators():
    separators = [
        "\n\n",
        "\n",
        " ",
        ".",
        ",",
        "\u200b",  # Zero-width space
        "\uff0c",  # Fullwidth comma
        "\u3001",  # Ideographic comma
        "\uff0e",  # Fullwidth full stop
        "\u3002",  # Ideographic full stop
        "",
    ]
    return separators


def load_pdf(pdf_path):
    """Load the pdf file."""
    doc = fitz.open(pdf_path)
    reader = easyocr.Reader(["en"], gpu=False)
    result = ""
    for i in range(doc.page_count):
        page = doc.load_page(i)
        pagetext = page.get_text().strip()
        if pagetext:
            if pagetext.endswith("!") or pagetext.endswith("?") or pagetext.endswith("."):
                result = result + pagetext
            else:
                result = result + pagetext + "."
        if len(doc.get_page_images(i)) > 0:
            for img in doc.get_page_images(i):
                if img:
                    pageimg = ""
                    xref = img[0]
                    img_data = doc.extract_image(xref)
                    img_bytes = img_data["image"]
                    pil_image = Image.open(io.BytesIO(img_bytes))
                    img = np.array(pil_image)
                    img_result = reader.readtext(img, paragraph=True, detail=0)
                    pageimg = pageimg + ", ".join(img_result).strip()
                    if pageimg.endswith("!") or pageimg.endswith("?") or pageimg.endswith("."):
                        pass
                    else:
                        pageimg = pageimg + "."
                result = result + pageimg
    return result


def load_html(html_path):
    """Load the html file."""
    data_html = UnstructuredHTMLLoader(html_path).load()
    content = ""
    for ins in data_html:
        content += ins.page_content
    return content


def load_txt(txt_path):
    """Load txt file."""
    with open(txt_path, "r") as file:
        text = file.read()
    return text


def load_doc(doc_path):
    """Load doc file."""
    logging.info("Converting doc file to docx file...")
    docx_path = doc_path + "x"
    os.system(f"libreoffice --headless --invisible --convert-to docx --outdir {os.path.dirname(docx_path)} {doc_path}")
    logging.info("Converted doc file to docx file.")
    text = load_docx(docx_path)
    os.remove(docx_path)
    return text


def load_docx(docx_path):
    """Load docx file."""
    doc = docx.Document(docx_path)
    text = ""
    # Save all 'rId:filenames' relationships in an dictionary and save the images if any.
    rid2img = {}
    for r in doc.part.rels.values():
        if isinstance(r._target, docx.parts.image.ImagePart):
            rid2img[r.rId] = os.path.basename(r._target.partname)
    if rid2img:
        save_path = "./imgs/"
        os.makedirs(save_path, exist_ok=True)
        docx2txt.process(docx_path, save_path)
    for paragraph in doc.paragraphs:
        if hasattr(paragraph, "text"):
            text += paragraph.text + "\n"
    if rid2img:
        shutil.rmtree(save_path)
    return text


def load_ppt(ppt_path):
    """Load ppt file."""
    logging.info("Converting ppt file to pptx file...")
    pptx_path = ppt_path + "x"
    os.system(f"libreoffice --headless --invisible --convert-to pptx --outdir {os.path.dirname(pptx_path)} {ppt_path}")
    logging.info("Converted ppt file to pptx file.")
    text = load_pptx(pptx_path)
    os.remove(pptx_path)
    return text


def load_pptx(pptx_path):
    """Load pptx file."""
    text = ""
    prs = pptx.Presentation(pptx_path)
    for slide in prs.slides:
        for shape in sorted(slide.shapes, key=lambda shape: (shape.top, shape.left)):
            if shape.has_text_frame:
                if shape.text:
                    text += shape.text + "\n"
            if shape.has_table:
                table_contents = "\n".join(
                    [
                        "\t".join([(cell.text if hasattr(cell, "text") else "") for cell in row.cells])
                        for row in shape.table.rows
                        if hasattr(row, "cells")
                    ]
                )
                if table_contents:
                    text += table_contents + "\n"

    return text


def load_md(md_path):
    """Load md file."""
    loader = UnstructuredMarkdownLoader(md_path)
    text = loader.load()[0].page_content
    return text


def load_xml(xml_path):
    """Load xml file."""
    loader = UnstructuredXMLLoader(xml_path)
    text = loader.load()[0].page_content
    return text


def load_json(json_path):
    """Load and process json file."""
    with open(json_path, "r") as file:
        data = json.load(file)
    return json.dumps(data)


def load_yaml(yaml_path):
    """Load and process yaml file."""
    with open(yaml_path, "r") as file:
        data = yaml.safe_load(file)
    return yaml.dump(data)


def load_xlsx(input_path):
    """Load and process xlsx file."""
    df = pd.read_excel(input_path)
    return df.to_string()


def load_csv(input_path):
    """Load the csv file."""
    df = pd.read_csv(input_path)
    return df.to_string()


def document_loader(doc_path):
    if doc_path.endswith(".pdf"):
        return load_pdf(doc_path)
    elif doc_path.endswith(".html"):
        return load_html(doc_path)
    elif doc_path.endswith(".txt"):
        return load_txt(doc_path)
    elif doc_path.endswith(".doc"):
        return load_doc(doc_path)
    elif doc_path.endswith(".docx"):
        return load_docx(doc_path)
    elif doc_path.endswith(".ppt"):
        return load_ppt(doc_path)
    elif doc_path.endswith(".pptx"):
        return load_pptx(doc_path)
    elif doc_path.endswith(".md"):
        return load_md(doc_path)
    elif doc_path.endswith(".xml"):
        return load_xml(doc_path)
    elif doc_path.endswith(".json") or doc_path.endswith(".jsonl"):
        return load_json(doc_path)
    elif doc_path.endswith(".yaml"):
        return load_yaml(doc_path)
    elif doc_path.endswith(".xlsx") or doc_path.endswith(".xls"):
        return load_xlsx(doc_path)
    elif doc_path.endswith(".csv"):
        return load_csv(doc_path)
    else:
        raise NotImplementedError(
            "Current only support pdf, html, txt, doc, docx, pptx, ppt, md, xml"
            + ", json, jsonl, yaml, xlsx, xls and csv format."
        )

def uni_pro(text):
    """Check if the character is ASCII or falls in the category of non-spacing marks."""
    normalized_text = unicodedata.normalize("NFKD", text)
    filtered_text = ""
    for char in normalized_text:
        if ord(char) < 128 or unicodedata.category(char) == "Mn":
            filtered_text += char
    return filtered_text


def load_html_data(url):
    crawler = Crawler()
    res = crawler.fetch(url)
    if res is None:
        return None
    soup = crawler.parse(res.text)
    all_text = crawler.clean_text(soup.select_one("body").text)
    main_content = ""
    for element_name in ["main", "container"]:
        main_block = None
        if soup.select(f".{element_name}"):
            main_block = soup.select(f".{element_name}")
        elif soup.select(f"#{element_name}"):
            main_block = soup.select(f"#{element_name}")
        if main_block:
            for element in main_block:
                text = crawler.clean_text(element.text)
                if text not in main_content:
                    main_content += f"\n{text}"
            main_content = crawler.clean_text(main_content)
    main_content = all_text if main_content == "" else main_content
    main_content = main_content.replace("\n", "")
    main_content = main_content.replace("\n\n", "")
    main_content = uni_pro(main_content)
    main_content = re.sub(r"\s+", " ", main_content)
    return main_content


def parse_html(input):
    """Parse the uploaded file."""
    chucks = []
    for link in input:
        if re.match(r"^https?:/{2}\w.+$", link):
            content = load_html_data(link)
            if content is None:
                continue
            chuck = [[content.strip(), link]]
            chucks += chuck
        else:
            logging.info("The given link/str {} cannot be parsed.".format(link))

    return chucks


def get_tables_result(pdf_path, table_strategy):
    """Extract tables information from pdf file."""
    if table_strategy == "fast":
        return None

    from unstructured.documents.elements import FigureCaption
    from unstructured.partition.pdf import partition_pdf

    tables_result = []
    raw_pdf_elements = partition_pdf(
        filename=pdf_path,
        infer_table_structure=True,
    )
    tables = [el for el in raw_pdf_elements if el.category == "Table"]
    for table in tables:
        table_coords = table.metadata.coordinates.points
        content = table.metadata.text_as_html
        table_page_number = table.metadata.page_number
        min_distance = float("inf")
        table_summary = None
        if table_strategy == "hq":
            for element in raw_pdf_elements:
                if isinstance(element, FigureCaption) or element.text.startswith("Tab"):
                    caption_page_number = element.metadata.page_number
                    caption_coords = element.metadata.coordinates.points
                    related, y_distance = get_relation(
                        table_coords, caption_coords, table_page_number, caption_page_number
                    )
                    if related:
                        if y_distance < min_distance:
                            min_distance = y_distance
                            table_summary = element.text
            if table_summary is None:
                parent_id = table.metadata.parent_id
                for element in raw_pdf_elements:
                    if element.id == parent_id:
                        table_summary = element.text
                        break
        elif table_strategy is None:
            table_summary = None
        if table_summary is None:
            text = f"[Table: {content}]"
        else:
            text = f"|Table: [Summary: {table_summary}], [Content: {content}]|"
        tables_result.append(text)
    return tables_result


def get_relation(table_coords, caption_coords, table_page_number, caption_page_number, threshold=100):
    """Get the relation of a pair of table and caption."""
    same_page = table_page_number == caption_page_number
    x_overlap = (min(table_coords[2][0], caption_coords[2][0]) - max(table_coords[0][0], caption_coords[0][0])) > 0
    if table_coords[0][1] - caption_coords[1][1] >= 0:
        y_distance = table_coords[0][1] - caption_coords[1][1]
    elif caption_coords[0][1] - table_coords[1][1] >= 0:
        y_distance = caption_coords[0][1] - table_coords[1][1]
    else:
        y_distance = 0
    y_close = y_distance < threshold
    return same_page and x_overlap and y_close, y_distance
