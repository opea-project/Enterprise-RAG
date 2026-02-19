# Copyright (C) 2024-2026 Intel Corporation
# SPDX-License-Identifier: Apache-2.0

import olefile
import os
import pptx
import re
import uuid
import zipfile

from pptx.exc import PackageNotFoundError

from comps.text_extractor.utils.file_loaders.abstract_loader import AbstractLoader
from comps.text_extractor.utils.file_loaders.load_image import LoadImage
from comps.cores.mega.logger import get_opea_logger, change_opea_logger_level

logger = get_opea_logger(f"{__file__.split('comps/')[1].split('/', 1)[0]}_microservice")
change_opea_logger_level(logger, log_level=os.getenv("OPEA_LOGGER_LEVEL", "INFO"))

class LoadPpt(AbstractLoader):
    """ Load and parse ppt/pptx files. """
    def is_encrypted(self):
        """ Check if the ppt file is encrypted. """
        ole = olefile.OleFileIO(self.file_path)
        if ole.exists('EncryptedPackage'):
            logger.info(f"[{self.file_path}] This file is encrypted.")
            return True
        return False

    def extract_text(self):
        """Extract text from PPT/PPTX (shapes, tables, notes, comments, and OCRed images)."""
        converted = False
        if self.file_type == "ppt":
            self.convert_to_pptx(self.file_path)
            converted = True

        try:
            text = ""
            prs = pptx.Presentation(self.file_path)
            logger.info(f"[{self.file_path}] Processing {len(prs.slides)} slides")
            
            # Extract comments from pptx file (stored separately in XML)
            comments_by_slide = self._extract_comments_from_xml()
            
            # Extract SmartArt text from diagram files
            smartart_texts = self._extract_graphic_text()

            # Iterate slides and collect slide-level header
            for idx, slide in enumerate(prs.slides, start=1):
                text += f"\n--- Slide {idx} ---\n"

                # Iterate shapes on slide (ordered by top/left)
                for shape in sorted(slide.shapes, key=lambda shape: (shape.top, shape.left)):
                    if shape.has_text_frame:
                        if shape.text:
                            text += shape.text + "\n"
                    if shape.has_table:
                        table_contents = "\n".join(
                            [
                                "\t".join([(cell.text if hasattr(cell, "text") else "") for cell in row.cells])
                                for row in shape.table.rows
                                if hasattr(row, "cells")
                            ]
                        )
                        if table_contents:
                            text += table_contents + "\n"
                            
                    # Image extraction + OCR: save temp image, run OCR, append when present
                    if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                        img_path = ""
                        try:
                            upload_path = os.getenv("UPLOAD_PATH", "/tmp/opea_upload")
                            os.makedirs(upload_path, exist_ok=True)
                            img_path = self.save_image(shape.image, upload_path)
                            if not img_path:
                                logger.error(f"[{self.file_path}] Failed to save image from slide {idx}")
                                continue
                            logger.debug(f"[{self.file_path}] Extracted {img_path} for processing")

                            ocr_text = LoadImage(img_path).extract_text().strip()
                            if ocr_text:
                                text += f"\n{ocr_text}\n"
                                logger.info(f"[{self.file_path}] OCRed image from slide {idx}: {len(ocr_text)} chars")
                                logger.debug(f"[{self.file_path}] OCRed image from slide {idx}: {ocr_text[:1000]}")
                            else:
                                logger.error(f"[{self.file_path}] OCR returned empty text for slide {idx}")
                        except Exception as e:
                            logger.error(f"[{self.file_path}] OCR failed on slide {idx}: {e}")
                        finally:
                            if img_path and os.path.isfile(img_path):
                                try:
                                    os.remove(img_path)
                                    logger.debug(f"[{self.file_path}] Removed {img_path}")
                                except Exception:
                                    logger.debug(f"[{self.file_path}] Failed to remove {img_path}")
                    
                    # Extract chart data (bar, pie, line, scatter, etc.)
                    if shape.has_chart:
                        try:
                            chart_text = self._extract_chart_data(shape.chart)
                            if chart_text:
                                text += f"\n{chart_text}\n"
                                logger.info(f"[{self.file_path}] Extracted chart from slide {idx}: {len(chart_text)} chars")
                                logger.debug(f"[{self.file_path}] Extracted chart from slide {idx}: {chart_text[:1000]}")
                        except Exception as e:
                            logger.error(f"[{self.file_path}] Error extracting chart from slide {idx}: {e}")
                    
                    # Extract SmartArt text content
                    if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                        try:
                            group_texts = self._extract_graphic_text(shape)
                            if group_texts:
                                group_text_joined = ' | '.join(group_texts)
                                text += f"[SmartArt/Group]: {group_text_joined}\n"
                                logger.info(f"[{self.file_path}] Extracted SmartArt/group text from slide {idx}: {len(group_text_joined)} chars")
                                logger.debug(f"[{self.file_path}] Extracted SmartArt/group text from slide {idx}: {group_text_joined[:1000]}")
                        except Exception as e:
                            logger.error(f"[{self.file_path}] Error extracting SmartArt from slide {idx}: {e}")
                
                # Extract speaker notes
                try:
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            text += f"[Slide Notes]: {notes_text}\n"
                            logger.debug(f"[{self.file_path}] Extracted notes from slide {idx}")
                except Exception as e:
                    logger.error(f"[{self.file_path}] Error extracting notes from slide {idx}: {e}. Continuing...")
                
                # Extract comments for this slide
                if idx in comments_by_slide:
                    for comment in comments_by_slide[idx]:
                        text += f"[Comment by {comment['author']}]: {comment['text']}\n"
                        logger.debug(f"[{self.file_path}] Extracted comment from slide {idx}")
            
            # Add SmartArt text at the end
            if smartart_texts:
                text += "\n[SmartArt Content]\n"
                for smartart_text in smartart_texts:
                    text += f"  {smartart_text}\n"
                logger.info(f"[{self.file_path}] Extracted {len(smartart_texts)} SmartArt text elements")
                    
        except PackageNotFoundError as e:
            if os.path.exists(self.file_path) and self.is_encrypted():
                logger.error(f"[{self.file_path}] The file is encrypted and cannot be processed: {e}")
                raise ValueError(f"The file {self.file_path} is encrypted and cannot be processed. Please decrypt it before uploading.")
            else:
                logger.error(f"[{self.file_path}] Error opening PPTX file: {e}")
                raise e
        finally:
            # Ensure the converted file is deleted. The original file deletion
            # should be handled in the caller code
            if converted and os.path.exists(self.file_path):
                os.remove(self.file_path)
                logger.info(f"Removed temporary converted file {self.file_path}")

        return text

    def _extract_comments_from_xml(self):
        """Extract comments from PPTX file by reading XML directly.
        Supports both legacy and modern formats.
        
        Returns:
            dict: Slide numbers mapped to lists of {'author', 'text'} dicts
        """
        comments_by_slide = {}
        
        try:
            with zipfile.ZipFile(self.file_path, 'r') as zip_ref:
                files = zip_ref.namelist()
                comment_files = [f for f in files if f.startswith('ppt/comments/') and f.endswith('.xml')]
                
                if not comment_files:
                    return comments_by_slide
                
                # Extract authors: legacy (numeric IDs) or modern (GUID IDs)
                authors_by_id = {}
                author_configs = [
                    ('ppt/commentAuthors.xml', r'<p:cmAuthor[^>]*(?:\sid=["\']?(\d+)["\']?[^>]*\sname=["\']([^"\']+)["\']|\sname=["\']([^"\']+)["\'][^>]*\sid=["\']?(\d+)["\']?)'),
                    ('ppt/authors.xml', r'<p\d+:author[^>]*(?:\sid=["\'](\{[^}]+\})["\'][^>]*\sname=["\']([^"\']+)["\']|\sname=["\']([^"\']+)["\'][^>]*\sid=["\'](\{[^}]+\})["\'])')
                ]
                
                for filename, pattern in author_configs:
                    if filename in files:
                        xml = zip_ref.read(filename).decode('utf-8')
                        for match in re.finditer(pattern, xml):
                            groups = match.groups()
                            author_id, author_name = (groups[0], groups[1]) if groups[0] else (groups[3], groups[2])
                            authors_by_id[author_id] = author_name
                
                # Extract comments by matching slide relationships
                for rel_file in [f for f in files if re.search(r'slides/_rels/slide(\d+)\.xml\.rels$', f)]:
                    slide_num = int(re.search(r'slide(\d+)', rel_file).group(1))
                    rels_content = zip_ref.read(rel_file).decode('utf-8')
                    
                    for comment_file in comment_files:
                        if comment_file.split('/')[-1] not in rels_content:
                            continue
                            
                        xml = zip_ref.read(comment_file).decode('utf-8')
                        # Unified pattern: <p[digits]:cm authorId="ID_OR_GUID"> or <p:cm authorId="ID">
                        for cm_match in re.finditer(r'<p\d*:cm[^>]*\sauthorId=["\']?([^"\'>\s]+)["\']?[^>]*>(.*?)</p\d*:cm>', xml, re.DOTALL):
                            author_id, cm_content = cm_match.groups()
                            author_name = authors_by_id.get(author_id, 'Unknown_author')
                            
                            for text_match in re.finditer(r'<a:t[^>]*>([^<]+)</a:t>', cm_content):
                                comment_text = text_match.group(1).strip()
                                if comment_text:
                                    comments_by_slide.setdefault(slide_num, []).append({
                                        'author': author_name,
                                        'text': comment_text
                                    })
                
                if comments_by_slide:
                    logger.info(f"[{self.file_path}] Extracted comments for slides: {list(comments_by_slide.keys())}")
                
        except Exception as e:
            logger.error(f"[{self.file_path}] Error extracting comments from XML: {e}")
        
        return comments_by_slide

    def _extract_chart_data(self, chart):
        """Extract data from chart including type, categories, series names and values."""
        from pptx.enum.chart import XL_CHART_TYPE

        # Map chart types to human-readable names
        chart_type_names = {
            XL_CHART_TYPE.BAR_CLUSTERED: "Bar Chart", XL_CHART_TYPE.BAR_STACKED: "Bar Chart",
            XL_CHART_TYPE.COLUMN_CLUSTERED: "Column Chart", XL_CHART_TYPE.COLUMN_STACKED: "Column Chart",
            XL_CHART_TYPE.LINE: "Line Chart", XL_CHART_TYPE.LINE_MARKERS: "Line Chart",
            XL_CHART_TYPE.PIE: "Pie Chart", XL_CHART_TYPE.PIE_EXPLODED: "Pie Chart",
            XL_CHART_TYPE.XY_SCATTER: "Scatter Plot", XL_CHART_TYPE.XY_SCATTER_LINES: "Scatter Plot",
            XL_CHART_TYPE.AREA: "Area Chart", XL_CHART_TYPE.DOUGHNUT: "Doughnut Chart",
        }

        # Extract chart title and build header with chart type
        title = chart.chart_title.text_frame.text.strip() if chart.has_title and chart.chart_title.has_text_frame else ""
        header = f"\n[{chart_type_names.get(chart.chart_type, f'Chart ({chart.chart_type})')}"
        text = f"{header + ': ' + title if title else header}]\n"

        # Extract category labels (x-axis) from first plot if available
        plot = chart.plots[0] if chart.plots else None
        categories = [str(cat) for cat in getattr(plot, "categories", []) or []]

        # Process each data series with its values and category mapping
        for series in chart.series:
            values = list(series.values or [])
            text += f"  {series.name or 'Series'}:\n"
            if categories and len(categories) == len(values):
                text += "".join(f"    {cat}: {val}\n" for cat, val in zip(categories, values))
            else:
                text += "".join(f"    Point {i+1}: {val}\n" for i, val in enumerate(values))

        return text

    def _extract_graphic_text(self, group_shape=None):
        """Extract text from SmartArt XML diagrams or grouped shapes.
        
        If group_shape is None, extracts from diagram XML files.
        Otherwise, recursively extracts from grouped shapes.
        Returns list of text entries.
        """
        texts = []
        
        if group_shape is None:
            # Extract text from SmartArt diagram XML files
            try:
                with zipfile.ZipFile(self.file_path, "r") as zf:
                    for name in zf.namelist():
                        if name.startswith("ppt/diagrams/data") and name.endswith(".xml"):
                            xml = zf.read(name).decode("utf-8")
                            for match in re.finditer(r"<a:t>([^<]+)</a:t>", xml):
                                entry = match.group(1).strip()
                                if entry and entry != "[Text]":
                                    texts.append(entry)
            except Exception as e:
                logger.error(f"[{self.file_path}] Error extracting SmartArt from diagram XML: {e}")
        else:
            # Recursively extract text from grouped shapes
            try:
                for shape in group_shape.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        texts.append(shape.text.strip())
                    if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                        nested = self._extract_graphic_text(shape)
                        texts.extend(nested)
            except Exception as e:
                logger.error(f"[{self.file_path}] Error extracting text from grouped shapes: {e}")
        
        return texts

    def save_image(self, image, save_path="/tmp/opea_upload"):
        if image:
            image_filename = os.path.join(save_path, f"{uuid.uuid4()}.{image.ext}")
            with open(image_filename, 'wb') as f:
                f.write(image.blob)
            return image_filename
        return None

    def convert_to_pptx(self, ppt_path):
        """Convert ppt file to pptx file."""
        pptx_path = ppt_path + "x"
        convert_log_file = f'/tmp/convert_{uuid.uuid4()}.log'
        exit_code = os.system(f"libreoffice --headless --invisible --convert-to pptx --outdir {os.path.dirname(pptx_path)} '{ppt_path}' > {convert_log_file} 2>&1")
        if exit_code != 0 or not os.path.exists(pptx_path):
            error = ""
            logger.error(f"Failed to convert {ppt_path} to pptx format. Exit code: {exit_code}")
            if os.path.exists(convert_log_file):
                try:
                    with open(convert_log_file, 'r') as f:
                        error = f.read()
                    logger.error(f"Conversion error: {error}")
                finally:
                    os.remove(convert_log_file)
            raise ValueError(f"Failed to convert {ppt_path} to pptx format. Error: {error}")
        else:
            self.file_path = pptx_path
            self.file_type = "pptx"
